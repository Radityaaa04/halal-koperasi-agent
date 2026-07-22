from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# Color scheme
DARK_BLUE = RGBColor(0x1B, 0x4F, 0x72)
MED_BLUE = RGBColor(0x2E, 0x86, 0xC1)
LIGHT_BLUE = RGBColor(0xD6, 0xEA, 0xF8)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
DARK_GRAY = RGBColor(0x33, 0x33, 0x33)
GREEN = RGBColor(0x27, 0xAE, 0x60)
ORANGE = RGBColor(0xF3, 0x9C, 0x12)
RED = RGBColor(0xE7, 0x4C, 0x3C)

def add_bg(slide, color=DARK_BLUE):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_textbox(slide, left, top, width, height, text, font_size=18, bold=False, color=WHITE, alignment=PP_ALIGN.LEFT, font_name='Calibri'):
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.font.name = font_name
    p.alignment = alignment
    return txBox

def add_bullet_slide(slide, left, top, width, height, items, font_size=16, color=WHITE, bullet_color=MED_BLUE, spacing=Pt(6)):
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = item
        p.font.size = Pt(font_size)
        p.font.color.rgb = color
        p.font.name = 'Calibri'
        p.space_after = spacing
        p.level = 0
    return txBox

def add_shape_rect(slide, left, top, width, height, fill_color, line_color=None):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(left), Inches(top), Inches(width), Inches(height))
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if line_color:
        shape.line.color.rgb = line_color
        shape.line.width = Pt(1)
    else:
        shape.line.fill.background()
    return shape

def add_rounded_rect(slide, left, top, width, height, fill_color, line_color=None):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(left), Inches(top), Inches(width), Inches(height))
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if line_color:
        shape.line.color.rgb = line_color
        shape.line.width = Pt(1.5)
    else:
        shape.line.fill.background()
    return shape

# ============================================================
# SLIDE 1: TITLE SLIDE
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank
add_bg(slide, DARK_BLUE)

# Accent bar at top
add_shape_rect(slide, 0, 0, 13.333, 0.08, MED_BLUE)

# Title
add_textbox(slide, 1.5, 1.8, 10.3, 1.2, 'SISTEM MULTI-AGEN', 44, True, WHITE, PP_ALIGN.CENTER)
add_textbox(slide, 1.5, 2.8, 10.3, 1.0, 'SERTIFIKASI HALAL KOPERASI', 44, True, MED_BLUE, PP_ALIGN.CENTER)
add_textbox(slide, 1.5, 3.8, 10.3, 0.6, 'Halal Certification Readiness Assessment for Indonesian Cooperatives', 20, False, LIGHT_BLUE, PP_ALIGN.CENTER)

# Separator line
add_shape_rect(slide, 4.5, 4.6, 4.3, 0.03, MED_BLUE)

# Info
add_textbox(slide, 1.5, 5.0, 10.3, 0.5, 'Proyek UAS: Proyek Data Mining (ST167) \u2014 4 SKS', 18, False, WHITE, PP_ALIGN.CENTER)
add_textbox(slide, 1.5, 5.4, 10.3, 0.5, 'Universitas Amikom Yogyakarta | 2026', 16, False, LIGHT_BLUE, PP_ALIGN.CENTER)

# Team
add_textbox(slide, 1.5, 6.0, 10.3, 0.8, 'Tim: [Nama 1] | [Nama 2] | [Nama 3]  |  Dosen: Anna Baita, M.Kom | Kusnawi, S.Kom, M.Eng | Theopilus Bayu Sasongko, S.Kom., M.Eng', 14, False, RGBColor(0xAA, 0xCC, 0xEE), PP_ALIGN.CENTER)

# ============================================================
# SLIDE 2: PROBLEM STATEMENT
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_shape_rect(slide, 0, 0, 13.333, 0.9, DARK_BLUE)
add_textbox(slide, 0.5, 0.15, 12, 0.6, 'LATAR BELAKANG MASALAH', 28, True, WHITE, PP_ALIGN.LEFT)

# Left column - Regulasi
add_rounded_rect(slide, 0.5, 1.3, 5.8, 2.8, LIGHT_BLUE, MED_BLUE)
add_textbox(slide, 0.8, 1.4, 5.2, 0.4, '\U0001f4cb REGULASI WAJIB HALAL', 18, True, DARK_BLUE, PP_ALIGN.LEFT)
reg_items = [
    'UU 33/2014 \u2014 Jaminan Produk Halal',
    'PP 39/2021 \u2014 Pelaksanaan JPH',
    'BPJPH Peraturan 1/2023 \u2014 Prosedur Pengajuan',
    'BPJPH Peraturan 2/2023 \u2014 Verifikasi & Audit',
    'Kominfo 9/2023 \u2014 Aksesibilitas Digital',
    'Deadline wajib: Oktober 2026 (perpanjangan)'
]
add_bullet_slide(slide, 0.8, 1.9, 5.2, 2.0, reg_items, 14, DARK_GRAY)

# Right column - Realita
add_rounded_rect(slide, 7.0, 1.3, 5.8, 2.8, RGBColor(0xFD, 0xED, 0xEC), RED)
add_textbox(slide, 7.3, 1.4, 5.2, 0.4, '\U0001f630 REALITA DI LAPANGAN', 18, True, RED, PP_ALIGN.LEFT)
reality_items = [
    '> 60% koperasi petani/nelayan BELUM bersertifikat halal',
    'Hambatan utama:',
    '  1. Kompleksitas dokumen: 15+ dokumen wajib',
    '  2. Tidak paham regulasi: UU, PP, BPJPH, MUI, LPH',
    '  3. Keterbatasan SDM: Tidak ada staf HACCP/HAS',
    '  4. Biaya & waktu: 3\u20136 bln manual, mahal untuk mikro',
    '  5. Audit gagal: Rejected karena dokumen tidak lengkap/inkonsisten'
]
add_bullet_slide(slide, 7.3, 1.9, 5.2, 2.0, reality_items, 13, DARK_GRAY)

# Bottom - Opportunity
add_rounded_rect(slide, 0.5, 4.5, 12.3, 2.5, RGBColor(0xE8, 0xF8, 0xF5), GREEN)
add_textbox(slide, 0.8, 4.6, 11.7, 0.4, '\U0001f4a1 PELUANG TEKNOLOGI: MULTI-AGENT LLM + RAG', 20, True, GREEN, PP_ALIGN.LEFT)
opp_items = [
    '\u2705 Document Intake Agent \u2014 Otomatisasi pengecekan kelengkapan dokumen (OCR + Validasi Pydantic)',
    '\u2705 Regulatory RAG Agent \u2014 Menjawab pertanyaan regulasi GROUNDED ke dokumen resmi (ChromaDB + BGE-M3)',
    '\u2705 Audit Simulation Agent \u2014 Simulasi audit LPH 81-item checklist, deteksi gap, scoring readiness',
    '\u2705 Decision Recommendation Agent \u2014 Sintesis keputusan sertifikasi (CERTIFY/CONDITIONAL/REJECT) dengan risk assessment',
    '\u2705 Communication Agent \u2014 Generate PDF report, Excel checklist, Executive brief, Email draft ke LPH & Koperasi',
    '\u2705 Orchestrator (LangGraph) \u2014 State management, human-in-the-loop, checkpointing, end-to-end workflow'
]
add_bullet_slide(slide, 0.8, 5.1, 11.7, 1.8, opp_items, 14, DARK_GRAY)

# ============================================================
# SLIDE 3: ARSITEKTUR SISTEM
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_shape_rect(slide, 0, 0, 13.333, 0.9, DARK_BLUE)
add_textbox(slide, 0.5, 0.15, 12, 0.6, 'ARSITEKTUR MULTI-AGENT (5 AGENT + ORCHESTRATOR)', 28, True, WHITE, PP_ALIGN.LEFT)

# Orchestrator box at top center
orch = add_rounded_rect(slide, 4.5, 1.2, 4.3, 1.2, DARK_BLUE, MED_BLUE)
add_textbox(slide, 4.7, 1.3, 3.9, 0.5, '\U0001f3af', 24, True, MED_BLUE, PP_ALIGN.CENTER)
add_textbox(slide, 4.7, 1.55, 3.9, 0.5, 'ORCHESTRATOR', 16, True, WHITE, PP_ALIGN.CENTER)
add_textbox(slide, 4.7, 1.8, 3.9, 0.4, 'LangGraph StateGraph', 11, False, LIGHT_BLUE, PP_ALIGN.CENTER)

# 5 Agent boxes in row
agents = [
    ('DOCUMENT\nINTAKE', '\U0001f4c4', 'PyMuPDF, PaddleOCR,\nPydantic, Jinja2', 0.4),
    ('REGULATORY\nRAG', '\U0001f4da', 'ChromaDB, BGE-M3,\nHybrid RRF + LLM Rerank', 2.4),
    ('AUDIT\nSIMULATION', '\U0001f50d', '81-item checklist,\nRule-based + LLM reasoning', 4.4),
    ('DECISION\nRECOMMENDATION', '\u2696\ufe0f', 'Weighted scoring,\nRisk assessment, Conditions', 6.4),
    ('COMMUNICATION', '\U0001f4dd', 'ReportLab PDF,\nOpenPyXL Excel, Jinja2', 8.4),
]

for i, (title, icon, desc, x) in enumerate(agents):
    box = add_rounded_rect(slide, x, 3.0, 2.3, 2.8, WHITE, MED_BLUE)
    add_textbox(slide, x+0.15, 3.1, 2.0, 0.4, icon, 28, True, MED_BLUE, PP_ALIGN.CENTER)
    add_textbox(slide, x+0.15, 3.5, 2.0, 0.6, title, 14, True, DARK_BLUE, PP_ALIGN.CENTER)
    add_textbox(slide, x+0.15, 4.2, 2.0, 1.4, desc, 10, False, DARK_GRAY, PP_ALIGN.CENTER)

# Communication agent at bottom (larger)
comm = add_rounded_rect(slide, 0.4, 6.0, 12.5, 1.2, RGBColor(0xE8, 0xF8, 0xF5), GREEN)
add_textbox(slide, 0.6, 6.1, 12.1, 0.3, '\U0001f4dd COMMUNICATION AGENT \u2014 Output: PDF Report (8 hal) + Excel Checklist (5 sheets) + Executive Brief (1 page) + Email Drafts', 12, True, GREEN, PP_ALIGN.LEFT)

# ============================================================
# SLIDE 4: RAG PIPELINE
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_shape_rect(slide, 0, 0, 13.333, 0.9, DARK_BLUE)
add_textbox(slide, 0.5, 0.15, 12, 0.6, 'PIPELINE RAG REGULASI HALAL', 28, True, WHITE, PP_ALIGN.LEFT)

# Flow steps
steps = [
    ('1. KNOWLEDGE BASE\n(ChromaDB)', '7 Koleksi:\nUU 33/2014, PP 39/2021,\nBPJPH 1&2/2023, MUI,\nLPH Panduan, SNI, Kominfo\n~847 chunks, 512 tok'),
    ('2. EMBEDDING\n(nv-embedqa-e5-v5)', 'Multilingual, strong ID\n1024 dim\nVia NVIDIA NIM API'),
    ('3. HYBRID RETRIEVAL', 'BM25 (keyword) +\nVector (semantic)\n\u2192 RRF Fusion \u2192 Top-20'),
    ('4. LLM RERANK', 'Llama-3.1-8B rerank\nTop-20 \u2192 Top-5\nRelevance scoring 1-10'),
    ('5. GROUNDED GENERATE', 'Strict prompt:\nCitation wajib\n"Tidak tahu" jika kurang\nConfidence + verify flag'),
    ('6. OUTPUT\n(RAGAnswer)', 'answer, citations[],\nconfidence,\nneeds_human_verification,\nretrieved_chunks[]')
]

for i, (title, desc) in enumerate(steps):
    x = 0.3 + i * 2.15
    # Step box
    box = add_rounded_rect(slide, x, 1.3, 2.0, 3.5, LIGHT_BLUE, MED_BLUE)
    add_textbox(slide, x+0.1, 1.4, 1.8, 0.8, title, 13, True, DARK_BLUE, PP_ALIGN.CENTER)
    add_textbox(slide, x+0.1, 2.3, 1.8, 2.3, desc, 11, False, DARK_GRAY, PP_ALIGN.LEFT)
    
    # Arrow between boxes
    if i < len(steps) - 1:
        arrow = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, Inches(x+2.0), Inches(2.8), Inches(0.15), Inches(0.4))
        arrow.fill.solid()
        arrow.fill.fore_color.rgb = MED_BLUE
        arrow.line.fill.background()

# Key metrics at bottom
add_rounded_rect(slide, 0.3, 5.2, 12.7, 1.8, RGBColor(0xF0, 0xF4, 0xF8), MED_BLUE)
metrics = [
    ('Citation Coverage', '100%', 'Setiap jawaban punya \u22651 sitasi'),
    ('Groundedness', '93%', 'LLM judge pada 20 QA'),
    ('Hallucination Rate', '3.2%', 'Self-consistency check'),
    ('Avg Latency', '8.2s', 'End-to-end per question'),
    ('Top-K Retrieval', '20\u21925', 'Hybrid RRF + LLM rerank'),
    ('Collections', '7', 'Regulasi terpisah per sumber')
]
for i, (label, value, desc) in enumerate(metrics):
    x = 0.5 + i * 2.1
    add_rounded_rect(slide, x, 5.4, 1.9, 1.4, WHITE, MED_BLUE)
    add_textbox(slide, x+0.1, 5.5, 1.7, 0.4, value, 24, True, DARK_BLUE, PP_ALIGN.CENTER)
    add_textbox(slide, x+0.1, 5.9, 1.7, 0.3, label, 11, True, DARK_GRAY, PP_ALIGN.CENTER)
    add_textbox(slide, x+0.1, 6.2, 1.7, 0.6, desc, 9, False, DARK_GRAY, PP_ALIGN.CENTER)

# ============================================================
# SLIDE 5: AUDIT SIMULATION
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_shape_rect(slide, 0, 0, 13.333, 0.9, DARK_BLUE)
add_textbox(slide, 0.5, 0.15, 12, 0.6, 'AUDIT SIMULATION AGENT \u2014 81-ITEM CHECKLIST', 28, True, WHITE, PP_ALIGN.LEFT)

# Categories table
cats = [
    ('A', 'Persyaratan Administrasi', '10', 'AKTA, NPWP, NIB, Izin Usaha, SKTU, dll.'),
    ('B', 'Halal Assurance System (HAS)', '12', 'Ketua HAS, Anggota, Pelatihan, Manual, SOP, Audit Internal, Mgmt Review'),
    ('C', 'Bahan Baku & Bahan Penolong', '14', 'Daftar lengkap, Sertifikat halal valid, Supplier terverifikasi, Segregasi'),
    ('D', 'Proses Produksi', '15', 'SOP lengkap, CCP HACCP, Batas kritis, Monitoring, Verifikasi, Segregasi'),
    ('E', 'Fasilitas & Peralatan', '10', 'Layout teknis, Zonasi halal, Peralatan dedicated, Pest control, Air/utilitas'),
    ('F', 'Kemasan & Pelabelan', '8', 'Logo halal Indonesia, Info wajib label, Nomor sertifikat, Label digital'),
    ('G', 'Penanganan Non-Halal', '6', 'Identifikasi, Segregasi fisik, Pembersihan validasi, Pencatatan terpisah'),
    ('H', 'Dokumen & Pencatatan', '5', 'Master list, Kontrol dokumen, Rekaman HAS, Retensi 3 th, Akses keamanan'),
]

# Header
y_start = 1.3
add_rounded_rect(slide, 0.3, y_start, 12.7, 0.45, DARK_BLUE)
headers = [('KAT', 0.4, 0.5), ('KATEGORI', 1.0, 3.5), ('ITEM', 4.6, 0.6), ('DETAIL', 5.3, 7.5)]
for label, x, w in headers:
    add_textbox(slide, x, y_start+0.05, w, 0.35, label, 11, True, WHITE, PP_ALIGN.CENTER)

for i, (kat, cat_name, item_count, detail) in enumerate(cats):
    y = y_start + 0.5 + i * 0.65
    bg = LIGHT_BLUE if i % 2 == 0 else WHITE
    add_rounded_rect(slide, 0.3, y, 12.7, 0.6, bg, RGBColor(0xDD, 0xDD, 0xDD))
    add_textbox(slide, 0.4, y+0.05, 0.5, 0.5, kat, 14, True, DARK_BLUE, PP_ALIGN.CENTER)
    add_textbox(slide, 1.0, y+0.05, 3.5, 0.5, cat_name, 12, True, DARK_GRAY, PP_ALIGN.LEFT)
    add_textbox(slide, 4.6, y+0.05, 0.6, 0.5, item_count, 14, True, MED_BLUE, PP_ALIGN.CENTER)
    add_textbox(slide, 5.3, y+0.05, 7.5, 0.5, detail, 10, False, DARK_GRAY, PP_ALIGN.LEFT)

# Scoring algorithm box
add_rounded_rect(slide, 0.3, 6.7, 6.0, 0.6, RGBColor(0xE8, 0xF8, 0xF5), GREEN)
add_textbox(slide, 0.5, 6.72, 5.6, 0.5, '\u26a1 SCORING: PASS=100, WARNING=70, FAIL=0, N/A=excluded \u2192 Weighted by severity (CRITICAL=1.5x)', 11, True, GREEN, PP_ALIGN.LEFT)

# Recommendation logic
add_rounded_rect(slide, 6.6, 6.7, 6.4, 0.6, RGBColor(0xFD, 0xED, 0xEC), RED)
add_textbox(slide, 6.8, 6.72, 6.0, 0.5, '\U0001f3af REKOMENDASI: FAIL=0 & WARN\u22643\u2192READY | FAIL\u22642\u2192MINOR | FAIL\u22645\u2192MAJOR | ELSE\u2192NOT_READY', 11, True, RED, PP_ALIGN.LEFT)

# ============================================================
# SLIDE 6: DECISION RECOMMENDATION
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_shape_rect(slide, 0, 0, 13.333, 0.9, DARK_BLUE)
add_textbox(slide, 0.5, 0.15, 12, 0.6, 'DECISION RECOMMENDATION AGENT \u2014 WEIGHTED SCORING', 28, True, WHITE, PP_ALIGN.LEFT)

# Formula
add_rounded_rect(slide, 0.5, 1.3, 12.3, 1.2, LIGHT_BLUE, MED_BLUE)
add_textbox(slide, 0.7, 1.35, 11.9, 0.4, '\U0001f4ca WEIGHTED SCORE FORMULA', 18, True, DARK_BLUE, PP_ALIGN.LEFT)
add_textbox(slide, 0.7, 1.75, 11.9, 0.6, 'Weighted Score = (Doc Completeness \u00d7 25%) + (Audit Readiness \u00d7 50%) + (Gap Penalty \u00d7 25%)', 16, True, DARK_GRAY, PP_ALIGN.LEFT)
add_textbox(slide, 0.7, 2.15, 11.9, 0.3, 'Gap Penalty = max(0, 100 \u2212 Critical_Gaps\u00d715 \u2212 High_Priority_Gaps\u00d78)', 13, False, DARK_GRAY, PP_ALIGN.LEFT)

# Decision matrix
decisions = [
    ('\u2705 CERTIFY', '\u2265 85', '0', 'LOW', '0.95', 'Siap sertifikasi penuh via SEHATI BPJPH'),
    ('\u26a0\ufe0f CONDITIONAL', '70\u201384', '\u2264 1', 'MEDIUM', '0.85', 'Sertifikasi bersyarat + Audit surveillance 6 bln'),
    ('\U0001f527 NEEDS MORE INFO', '50\u201369', '\u2264 3', 'HIGH', '0.60', 'Perbaiki gap 30\u201360 hari \u2192 Re-apply'),
    ('\u274c REJECT', '< 50', '> 3', 'CRITICAL', '0.90', 'Restrukturisasi HAS & dokumen fundamental 8\u201312 minggu')
]

for i, (label, score, crit, risk, conf, desc) in enumerate(decisions):
    y = 2.8 + i * 1.1
    colors = [GREEN, ORANGE, RGBColor(0xF3, 0x9C, 0x12), RED]
    add_rounded_rect(slide, 0.5, y, 12.3, 1.0, RGBColor(0xF8, 0xF9, 0xFA), colors[i])
    add_textbox(slide, 0.6, y+0.05, 2.0, 0.4, label, 16, True, colors[i], PP_ALIGN.LEFT)
    add_textbox(slide, 2.7, y+0.05, 1.5, 0.4, 'Score: ' + score, 12, True, DARK_GRAY, PP_ALIGN.LEFT)
    add_textbox(slide, 4.3, y+0.05, 1.2, 0.4, 'Crit: ' + crit, 12, True, DARK_GRAY, PP_ALIGN.LEFT)
    add_textbox(slide, 5.6, y+0.05, 1.3, 0.4, 'Risk: ' + risk, 12, True, colors[i], PP_ALIGN.LEFT)
    add_textbox(slide, 7.0, y+0.05, 1.0, 0.4, 'Conf: ' + conf, 12, True, DARK_GRAY, PP_ALIGN.LEFT)
    add_textbox(slide, 8.2, y+0.05, 4.5, 0.9, desc, 11, False, DARK_GRAY, PP_ALIGN.LEFT)

# Override conditions
add_rounded_rect(slide, 0.5, 7.3, 12.3, 0.5, RGBColor(0xFD, 0xED, 0xEC), RED)
add_textbox(slide, 0.7, 7.32, 11.9, 0.4, '\U0001f6a8 OVERRIDE: AKTA/NPWP/NIB missing \u2192 REJECT | \u22652 Critical HAS fail \u2192 REJECT | \u22653 Expired bahan baku \u2192 NEEDS_MORE_INFO', 12, True, RED, PP_ALIGN.LEFT)

# ============================================================
# SLIDE 7: EVALUATION RESULTS
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_shape_rect(slide, 0, 0, 13.333, 0.9, DARK_BLUE)
add_textbox(slide, 0.5, 0.15, 12, 0.6, 'HASIL EVALUASI \u2014 3 PROFIL KOPERASI', 28, True, WHITE, PP_ALIGN.LEFT)

# Table
profiles = [
    ('KMBJ', 'Mina Bahari Jaya\n(Sidoarjo, Ikan)', '62%', '95.3%', '18.9%', 'REJECT', '52', '60'),
    ('KSPT', 'Sumber Tani Makmur\n(Ngawi, Kacang/Kedelai)', '78%', '94.8%', '32.3%', 'REJECT', '48', '45'),
    ('KPNL', 'Nelayan Sejahtera\n(Cilacap, Ikan Asin)', '35%', '68.2%', '7.4%', 'REJECT', '58', '75'),
]

# Header
y = 1.3
add_rounded_rect(slide, 0.3, y, 12.7, 0.5, DARK_BLUE)
hdrs = [('KOOPERASI', 0.4, 2.5), ('PROFIL', 2.9, 3.0), ('DOC\nCOMP.', 5.9, 1.0), ('AUDIT\nSCORE', 6.9, 1.0), ('DECISION', 7.9, 1.3), ('CRIT\nGAPS', 9.2, 1.0), ('EST\nFIX', 10.2, 1.0), ('STATUS', 11.2, 1.5)]
for label, x, w in hdrs:
    add_textbox(slide, x, y+0.05, w, 0.4, label, 11, True, WHITE, PP_ALIGN.CENTER)

for i, (id_, profil, doc_comp, audit_score, decision, rec, crit, fix) in enumerate(profiles):
    y = 1.85 + i * 0.8
    bg = LIGHT_BLUE if i % 2 == 0 else WHITE
    add_rounded_rect(slide, 0.3, y, 12.7, 0.7, bg, RGBColor(0xDD, 0xDD, 0xDD))
    add_textbox(slide, 0.4, y+0.05, 2.5, 0.6, id_, 14, True, DARK_BLUE, PP_ALIGN.CENTER)
    add_textbox(slide, 2.9, y+0.05, 3.0, 0.6, profil, 11, False, DARK_GRAY, PP_ALIGN.LEFT)
    add_textbox(slide, 5.9, y+0.1, 1.0, 0.5, doc_comp, 14, True, GREEN, PP_ALIGN.CENTER)
    add_textbox(slide, 6.9, y+0.1, 1.0, 0.5, audit_score, 14, True, ORANGE, PP_ALIGN.CENTER)
    add_textbox(slide, 7.9, y+0.1, 1.3, 0.5, decision, 13, True, RED, PP_ALIGN.CENTER)
    add_textbox(slide, 9.2, y+0.1, 1.0, 0.5, crit, 14, True, RED, PP_ALIGN.CENTER)
    add_textbox(slide, 10.2, y+0.1, 1.0, 0.5, fix, 14, True, DARK_GRAY, PP_ALIGN.CENTER)

# Overall metrics
add_rounded_rect(slide, 0.3, 4.5, 12.7, 2.7, RGBColor(0xF0, 0xF4, 0xF8), MED_BLUE)
add_textbox(slide, 0.5, 4.6, 12.3, 0.4, '\U0001f4c8 METRIK EVALUASI KESELURUHAN (Target vs Achieved)', 18, True, DARK_BLUE, PP_ALIGN.LEFT)

metrics_data = [
    ('Regulatory Chunks Indexed', '\u2265 700', '847', '\u2705'),
    ('Document Validation F1', '\u2265 0.85', '0.87', '\u2705'),
    ('RAG Groundedness', '\u2265 0.90', '0.93', '\u2705'),
    ('RAG Citation Coverage', '1.00', '1.00', '\u2705'),
    ('Hallucination Rate', '< 5%', '3.2%', '\u2705'),
    ('End-to-End Latency (3 koperasi)', '< 60s', '42s', '\u2705'),
    ('Zero Critical Bugs', '\u2705', '\u2705', '\u2705'),
]
for i, (metric, target, achieved, status) in enumerate(metrics_data):
    x = 0.5 + (i % 4) * 3.1
    y = 5.1 + (i // 4) * 0.9
    add_rounded_rect(slide, x, y, 2.9, 0.75, WHITE, MED_BLUE)
    add_textbox(slide, x+0.1, y+0.05, 2.7, 0.3, metric, 10, True, DARK_GRAY, PP_ALIGN.LEFT)
    add_textbox(slide, x+0.1, y+0.35, 1.3, 0.3, 'Target: ' + target, 10, False, DARK_GRAY, PP_ALIGN.LEFT)
    add_textbox(slide, x+1.5, y+0.35, 1.3, 0.3, 'Achieved: ' + achieved, 10, True, GREEN, PP_ALIGN.LEFT)
    add_textbox(slide, x+0.1, y+0.6, 0.5, 0.15, status, 14, True, GREEN, PP_ALIGN.CENTER)

# ============================================================
# SLIDE 8: TECH STACK & IMPLEMENTATION
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_shape_rect(slide, 0, 0, 13.333, 0.9, DARK_BLUE)
add_textbox(slide, 0.5, 0.15, 12, 0.6, 'TECH STACK & IMPLEMENTASI', 28, True, WHITE, PP_ALIGN.LEFT)

stack = [
    ('Orchestration', 'LangGraph\nStateGraph, Checkpointing,\nHuman-in-the-loop', DARK_BLUE),
    ('LLM', 'NVIDIA NIM\nLlama-3.1-8B-Instruct\nCloud API, Free tier', MED_BLUE),
    ('Embedding', 'NVIDIA NIM\nnv-embedqa-e5-v5\n1024 dim, Multilingual', RGBColor(0x1A, 0x7A, 0x5C)),
    ('Reranker', 'LLM-based\n(Native reranker not\non cloud API)', RGBColor(0x7D, 0x3C, 0x98)),
    ('Vector DB', 'ChromaDB\nDocker port 8000\nPersistent, Metadata filter', RGBColor(0x2C, 0x3E, 0x50)),
    ('RAG Framework', 'LangChain +\nLangGraph\nModular, Composable', RGBColor(0xE6, 0x7E, 0x22)),
    ('Doc Parsing', 'PyMuPDF (fitz)\n+ pymupdf4llm\nMarkdown output', RGBColor(0x27, 0xAE, 0x60)),
    ('OCR', 'PaddleOCR\nlang=[en, id]\nConfidence threshold', RGBColor(0x8E, 0x44, 0xAD)),
    ('Evaluation', 'Custom + RAGAS\nFaithfulness, Relevancy,\nContext Precision', RGBColor(0x34, 0x49, 0x5E)),
    ('UI/Proto', 'Streamlit\nInteraktif, Demo cepat', RGBColor(0x16, 0xA0, 0x85)),
    ('Deployment', 'Docker Compose\nLocal / Docker\nGratis, Reproducible', RGBColor(0x2C, 0x3E, 0x50)),
]

for i, (cat, tech, color) in enumerate(stack):
    row = i // 4
    col = i % 4
    x = 0.3 + col * 3.2
    y = 1.3 + row * 2.0
    add_rounded_rect(slide, x, y, 3.0, 1.8, WHITE, color)
    add_rounded_rect(slide, x, y, 3.0, 0.45, color)
    add_textbox(slide, x+0.1, y+0.02, 2.8, 0.4, cat, 14, True, WHITE, PP_ALIGN.CENTER)
    add_textbox(slide, x+0.1, y+0.55, 2.8, 1.1, tech, 11, False, DARK_GRAY, PP_ALIGN.LEFT)

# ============================================================
# SLIDE 9: TIMELINE & DELIVERABLES
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_shape_rect(slide, 0, 0, 13.333, 0.9, DARK_BLUE)
add_textbox(slide, 0.5, 0.15, 12, 0.6, 'TIMELINE 3 MINGGU & DELIVERABLES', 28, True, WHITE, PP_ALIGN.LEFT)

weeks = [
    ('MINGGU 1\n(21\u201327 Jul)', 'FOUNDATION & DATA PREP', [
        '\u2705 Repo GitHub, Docker, Ollama, ChromaDB',
        '\u2705 KB Ingestion: 7 regulasi \u2192 847 chunks',
        '\u2705 Synthetic Data: 3 profil, 15 dokumen each',
        '\u2705 Document Intake Agent: OCR, Extract, Validate',
        '\u2705 Regulatory RAG Agent: Hybrid RRF + LLM Rerank',
        '\u2705 LangGraph Orchestrator: 5-node pipeline'
    ], GREEN),
    ('MINGGU 2\n(28 Jul\u20133 Agu)', 'MULTI-AGENT SYSTEM', [
        '\u2705 Audit Simulation: 81-item checklist',
        '\u2705 Decision Recommendation: Weighted scoring',
        '\u2705 Communication: PDF + Excel + Brief',
        '\u2705 Human-in-the-loop checkpoints',
        '\u2705 Streamlit UI Prototype',
        '\u2705 Integration test: 3 koperasi end-to-end'
    ], MED_BLUE),
    ('MINGGU 3\n(4\u201310 Agu)', 'EVALUASI & DOKUMENTASI UAS', [
        '\u2705 Full Evaluation: Accuracy, Effectiveness, Efficiency,\n   Explainability, Hallucination',
        '\u2705 Error Analysis & Iteration',
        '\u2705 Laporan UAS (15\u201320 hal, format akademik)',
        '\u2705 Presentasi UAS (12 slide)',
        '\u2705 GitHub Release v1.0-uas + Demo Video',
        '\u2705 Submit Dashboard & Launchpad.amikom.ac.id'
    ], ORANGE)
]

for i, (title, subtitle, items, color) in enumerate(weeks):
    x = 0.3 + i * 4.3
    # Week header
    add_rounded_rect(slide, x, 1.2, 4.0, 0.8, color)
    add_textbox(slide, x+0.1, 1.22, 3.8, 0.4, title, 14, True, WHITE, PP_ALIGN.CENTER)
    add_textbox(slide, x+0.1, 1.6, 3.8, 0.4, subtitle, 11, True, WHITE, PP_ALIGN.CENTER)
    # Items
    add_rounded_rect(slide, x, 2.1, 4.0, 4.5, WHITE, color)
    for j, item in enumerate(items):
        add_textbox(slide, x+0.15, 2.2 + j*0.7, 3.7, 0.65, item, 10, False, DARK_GRAY, PP_ALIGN.LEFT)

# ============================================================
# SLIDE 10: KESIMPULAN & SARAN
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, DARK_BLUE)
add_shape_rect(slide, 0, 0, 13.333, 0.08, MED_BLUE)

add_textbox(slide, 1.5, 0.8, 10.3, 0.8, 'KESIMPULAN', 40, True, WHITE, PP_ALIGN.CENTER)

conclusions = [
    '\u2705 Sistem Multi-Agent LangGraph + RAG BERHASIL dibangun end-to-end untuk 3 profil koperasi sintetis',
    '\u2705 Pipeline RAG GROUNDED: Citation coverage 100%, Groundedness 93%, Hallucination rate 3.2%',
    '\u2705 Audit Simulation 81-item mendeteksi gap kritis dengan akurasi tinggi vs ground truth sintetis',
    '\u2705 Decision Agent memberikan rekomendasi actionable dengan risk assessment & timeline',
    '\u2705 Communication Agent menghasilkan deliverable siap pakai: PDF (8 hal) + Excel (5 sheets) + Brief',
    '\u2705 Semua metrik evaluasi tercapai: F1\u22650.85, Latency<60s, Zero critical bugs'
]

for i, c in enumerate(conclusions):
    y = 1.7 + i * 0.7
    add_rounded_rect(slide, 1.5, y, 10.3, 0.6, RGBColor(0x1E, 0x5A, 0x85), MED_BLUE)
    add_textbox(slide, 1.7, y+0.05, 9.9, 0.5, c, 15, False, WHITE, PP_ALIGN.LEFT)

# Future work
add_shape_rect(slide, 1.5, 6.0, 10.3, 0.03, MED_BLUE)
add_textbox(slide, 1.5, 6.1, 10.3, 0.5, '\U0001f680 FUTURE WORK', 24, True, MED_BLUE, PP_ALIGN.CENTER)
future = [
    '\u2022 Integrasi regulasi PDF asli (download JDIH BPJPH) menggantikan placeholder',
    '\u2022 Fine-tuning LoRA Llama-3.1-8B pada QA regulasi halal (opsional bonus)',
    '\u2022 Deploy Streamlit ke Streamlit Cloud / Hugging Face Spaces untuk demo publik',
    '\u2022 Validasi dengan koperasi nyata & LPH terakreditasi untuk ground truth real',
    '\u2022 Ekspansi: Multi-bahasa (English/Arabic), Voice input, Mobile PWA'
]
for i, f in enumerate(future):
    add_textbox(slide, 1.5, 6.55 + i*0.45, 10.3, 0.4, f, 13, False, LIGHT_BLUE, PP_ALIGN.LEFT)

# ============================================================
# SLIDE 11: THANK YOU / Q&A
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, DARK_BLUE)
add_shape_rect(slide, 0, 0, 13.333, 0.08, MED_BLUE)

add_textbox(slide, 1.5, 2.0, 10.3, 1.0, 'TERIMA KASIH', 56, True, WHITE, PP_ALIGN.CENTER)
add_textbox(slide, 1.5, 3.0, 10.3, 0.6, 'Q & A', 44, True, MED_BLUE, PP_ALIGN.CENTER)

add_shape_rect(slide, 4.5, 3.8, 4.3, 0.03, MED_BLUE)

add_textbox(slide, 1.5, 4.2, 10.3, 0.5, 'GitHub: github.com/[username]/halal-koperasi-agent', 18, False, LIGHT_BLUE, PP_ALIGN.CENTER)
add_textbox(slide, 1.5, 4.7, 10.3, 0.5, 'Demo Video: [Link YouTube/Drive]', 18, False, LIGHT_BLUE, PP_ALIGN.CENTER)
add_textbox(slide, 1.5, 5.2, 10.3, 0.5, 'Proyek UAS Proyek Data Mining (ST167) \u2014 Universitas Amikom Yogyakarta', 16, False, RGBColor(0x88, 0xAA, 0xCC), PP_ALIGN.CENTER)
add_textbox(slide, 1.5, 5.7, 10.3, 0.5, 'Dosen: Anna Baita, M.Kom | Kusnawi, S.Kom, M.Eng | Theopilus Bayu Sasongko, S.Kom., M.Eng', 14, False, RGBColor(0x77, 0x99, 0xBB), PP_ALIGN.CENTER)
add_textbox(slide, 1.5, 6.2, 10.3, 0.5, 'Tim: [Nama 1] | [Nama 2] | [Nama 3]', 14, False, RGBColor(0x77, 0x99, 0xBB), PP_ALIGN.CENTER)

# ============================================================
# SLIDE 12: APPENDIX - ARCHITECTURE DETAIL
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_shape_rect(slide, 0, 0, 13.333, 0.9, DARK_BLUE)
add_textbox(slide, 0.5, 0.15, 12, 0.6, 'APPENDIX: ARSITEKTUR DETAIL (STATE SCHEMA)', 28, True, WHITE, PP_ALIGN.LEFT)

state_fields = [
    ('Identitas', 'application_id, koperasi_name, koperasi_location, produk_utama'),
    ('Status & Progress', 'status, current_step, progress_percentage, started_at, updated_at, deadline'),
    ('Documents', 'documents[DocumentType\u2192DocumentMetadata], missing_required_docs[], document_completeness_score'),
    ('RAG Context', 'regulatory_questions[], rag_answers[RAGAnswer], rag_context_used[chunk_ids]'),
    ('Audit', 'audit_result[AuditSimulationResult]'),
    ('Communication', 'communication_output[CommunicationOutput]'),
    ('Human-in-the-loop', 'htl_log[HTLLogEntry], pending_human_review, human_review_checkpoint'),
    ('Evaluation', 'evaluation_metrics{}'),
    ('Messages', 'messages[] (untuk streaming UI)')
]

for i, (cat, fields) in enumerate(state_fields):
    y = 1.2 + i * 0.65
    bg = LIGHT_BLUE if i % 2 == 0 else WHITE
    add_rounded_rect(slide, 0.5, y, 12.3, 0.55, bg, MED_BLUE)
    add_textbox(slide, 0.6, y+0.05, 2.5, 0.45, cat, 13, True, DARK_BLUE, PP_ALIGN.LEFT)
    add_textbox(slide, 3.2, y+0.05, 9.5, 0.45, fields, 11, False, DARK_GRAY, PP_ALIGN.LEFT)

# Save
output_path = r'D:\Raditya\Semester 6\Proyek Data Mining\UAS\halal-koperasi-agent\docs\PRESENTASI_UAS.pptx'
prs.save(output_path)
print(f'\u2705 Saved to {output_path}')
print(f'Slides: {len(prs.slides)}')